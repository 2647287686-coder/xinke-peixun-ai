# -*- coding: utf-8 -*-
"""
kb/sync.py —— 把 workbuddy「新员工培训」资料库同步到本地知识库（kb/raw + kb/text）

特点：
  - 自包含：不依赖 WorkBuddy 插件缓存，标准 requests 调用，可在任意环境运行。
  - 动态列举：递归遍历空间所有节点，不再写死节点列表（新增资料也能抓到）。
  - 自动抽取：下载的 docx/xlsx/pptx/txt 直接转成 kb/text 下的纯文本片段源。

用法：
  # 通过环境变量注入当前 workbuddy token（token 由 WorkBuddy 会话换取，约 30 分钟有效）
  KB_TOKEN=op_xxx python kb/sync.py
  # 或命令行参数
  python kb/sync.py --token op_xxx

说明：
  - 资料库 API 的 token 是「会话级、约 30 分钟过期」，因此无法在公网后端无人值守地
    长期自动同步。正确做法是：在 WorkBuddy 会话内（能换取有效 token）运行本脚本，
    把最新资料同步到 kb/text，再 git push，Render 会自动重新部署，员工即看到最新知识库。
  - 也可在 app.py 运行的主机上用环境变量 KB_TOKEN 提供 token，并访问 /api/sync 触发同步；
    但务必注意 token 过期后需重新注入。
"""
import os, re, sys, argparse
import requests

API_BASE = os.environ.get("KB_API_BASE", "https://www.workbuddy.cn")
SPACE_ID = os.environ.get("KB_SPACE_ID", "obj0gboWcbKCS6OI3nrxWj")
HERE = os.path.dirname(os.path.abspath(__file__))
RAW = os.path.join(HERE, "raw")
TEXT = os.path.join(HERE, "text")
os.makedirs(RAW, exist_ok=True)
os.makedirs(TEXT, exist_ok=True)


def _headers(token):
    return {"X-Skill-Token": token, "Content-Type": "application/json"}


def list_nodes(token, space_id, parent_id=""):
    r = requests.post(f"{API_BASE}/space/api/agent/v1/list-node",
                      json={"spaceId": space_id, "parentNodeId": parent_id},
                      headers=_headers(token), timeout=20)
    r.raise_for_status()
    d = r.json()
    if str(d.get("code")) not in ("0", "OK", "ok", "None", ""):
        raise RuntimeError(f"list-node 失败: {d.get('code')} {d.get('msg')}")
    return d.get("data", {}).get("nodes", []) or []


def get_drive_file(token, node_id):
    r = requests.post(f"{API_BASE}/space/api/agent/v1/get-drive-file",
                      json={"nodeId": node_id},
                      headers=_headers(token), timeout=20)
    r.raise_for_status()
    d = r.json()
    data = d.get("data", {}) or {}
    return data.get("downloadUrl"), data.get("fileName"), data.get("ext")


def collect_files(token, space_id, parent_id="", depth=0):
    """递归收集所有可下载的文件节点（kind=drive 且有 url / 非文件夹）。"""
    out = []
    for n in list_nodes(token, space_id, parent_id):
        nid = n.get("id")
        if not nid:
            continue
        children = n.get("nodes") or []
        if children:
            out += collect_files(token, space_id, nid, depth + 1)
            continue
        kind = n.get("kind")
        if kind in ("drive",) or n.get("url"):
            out.append((nid, n.get("title") or nid))
    return out


def sanitize(name):
    return re.sub(r'[\\/:*?"<>|]', '_', name).strip()

def norm_text_name(title):
    # 仅保留中文/字母/数字，避免标题中空格、标点差异造成重复文本文件（覆盖写而非残留）
    s = re.sub(r'[^\u4e00-\u9fff0-9a-zA-Z]+', '', title)
    return s or sanitize(title)


# ---- 文本抽取（与 extract_text.py 一致，内联以便自包含）----
def extract_docx(path):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)

def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"=== 表: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v).strip() for v in row]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- 幻灯片 {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    line = " | ".join(c for c in cells if c)
                    if line:
                        parts.append(line)
    return "\n".join(parts)

def extract_txt(path):
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()


def download_and_extract(token):
    files = collect_files(token, SPACE_ID)
    print(f"[SYNC] 发现 {len(files)} 个文件节点")
    ok = fail = 0
    for nid, title in files:
        try:
            url, fname, ext = get_drive_file(token, nid)
            if not url:
                print(f"[SKIP] {title}: 无下载链接")
                continue
            ext = ext or (os.path.splitext(fname)[1].lstrip(".") if fname else "")
            ext = ext.lower()
            out_name = f"{sanitize(title)}.{ext}" if ext else sanitize(title)
            r = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=60)
            r.raise_for_status()
            raw_path = os.path.join(RAW, out_name)
            with open(raw_path, "wb") as f:
                f.write(r.content)
            # 抽取
            if ext == "docx":
                text = extract_docx(raw_path)
            elif ext == "xlsx":
                text = extract_xlsx(raw_path)
            elif ext == "pptx":
                text = extract_pptx(raw_path)
            elif ext == "txt":
                text = extract_txt(raw_path)
            else:
                print(f"[SKIP] {title}.{ext}: 非文本格式，跳过抽取")
                continue
            text = text.strip()
            if not text:
                print(f"[EMPTY] {title}.{ext}")
                continue
            with open(os.path.join(TEXT, norm_text_name(title) + ".txt"), "w", encoding="utf-8") as f:
                f.write(text)
            print(f"[OK] {title}.{ext} -> {len(text)} 字")
            ok += 1
        except Exception as e:
            print(f"[ERR] {title}: {e}")
            fail += 1
        time.sleep(0.2)
    print(f"\n[SYNC] 完成 抽取 {ok} 篇, 失败 {fail} 篇")
    return ok, fail


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--token", default=os.environ.get("KB_TOKEN", ""))
    args = ap.parse_args()
    if not args.token:
        print("错误：请提供 token（--token 或环境变量 KB_TOKEN）。\n"
              "token 由 WorkBuddy 会话换取，约 30 分钟有效。")
        sys.exit(2)
    try:
        ok, fail = download_and_extract(args.token)
        sys.exit(1 if fail and not ok else 0)
    except Exception as e:
        print(f"[SYNC] 同步失败: {e}")
        sys.exit(1)


if __name__ == "__main__":
    import time
    main()
