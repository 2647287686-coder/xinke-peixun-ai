import os, glob, sys

VENV = "C:/Users/admin/.workbuddy/binaries/python/envs/default/Scripts/python.exe"
RAW = "C:/Users/admin/WorkBuddy/2026-08-14-09-03-20/kb/raw"
OUT = "C:/Users/admin/WorkBuddy/2026-08-14-09-03-20/kb/text"
os.makedirs(OUT, exist_ok=True)

def extract_docx(path):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)

def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"=== 表: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v).strip() for v in row]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- 幻灯片 {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    line = " | ".join(c for c in cells if c)
                    if line:
                        parts.append(line)
    return "\n".join(parts)

def extract_txt(path):
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()

def main():
    count = 0
    for path in glob.glob(os.path.join(RAW, "*")):
        ext = os.path.splitext(path)[1].lower()
        name = os.path.splitext(os.path.basename(path))[0]
        try:
            if ext == ".docx":
                text = extract_docx(path)
            elif ext == ".xlsx":
                text = extract_xlsx(path)
            elif ext == ".pptx":
                text = extract_pptx(path)
            elif ext == ".txt":
                text = extract_txt(path)
            else:
                continue  # 图片/视频跳过
            text = text.strip()
            if not text:
                print(f"[EMPTY] {name}{ext}")
                continue
            with open(os.path.join(OUT, name + ".txt"), "w", encoding="utf-8") as f:
                f.write(text)
            print(f"[OK] {name}{ext} -> {len(text)} chars")
            count += 1
        except Exception as e:
            print(f"[ERR] {name}{ext}: {e}")
    print(f"\nExtracted {count} text documents to {OUT}")

if __name__ == "__main__":
    main()
